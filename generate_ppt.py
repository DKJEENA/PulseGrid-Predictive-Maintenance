from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    # Create presentation
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Predictive Maintenance Platform"
    subtitle.text = "AI-Driven IIoT Solution\nCollege Project Presentation"

    # Problem Statement Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Problem Statement"
    tf = body_shape.text_frame
    tf.text = "Industries lose massive revenue due to unplanned downtime."
    p = tf.add_paragraph()
    p.text = "Machines fail unexpectedly, leading to halted production."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Traditional maintenance is either too early (wasting parts) or too late (causing damage)."
    p.level = 1

    # Solution Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Our Solution"
    tf = body_shape.text_frame
    tf.text = "An end-to-end Predictive Maintenance Platform."
    p = tf.add_paragraph()
    p.text = "Ingests real-time sensor data (Vibration, Temp, RPM)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Uses Machine Learning (Random Forest) to predict failure risk."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Provides a real-time tracking dashboard for maintenance teams."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Features an Automated Dataset Auto-Cleaner for context-swapping."
    p.level = 1

    # Architecture Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "System Architecture"
    tf = body_shape.text_frame
    tf.text = "Technology Stack:"
    p = tf.add_paragraph()
    p.text = "Frontend: React + Vite + Recharts (Modern Dark UI)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend: FastAPI (Python REST API)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Machine Learning: Scikit-Learn (Joblib pipeline)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database: SQLite/PostgreSQL"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Data Source: AI4I 2020 Predictive Maintenance Dataset"
    p.level = 1

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "Predictive_Maintenance_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
