import os
import sys
from datetime import datetime

# Allow loading packages installed into workspace-local vendor folder
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
VENDOR = os.path.join(ROOT, ".vendor")
if VENDOR not in sys.path:
    sys.path.insert(0, VENDOR)

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = os.path.join(ROOT, "deliverables")
os.makedirs(OUT_DIR, exist_ok=True)
OUT_FILE = os.path.join(OUT_DIR, "PulseGrid_Project_Presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

COLOR_BG = RGBColor(9, 23, 39)
COLOR_PANEL = RGBColor(20, 48, 71)
COLOR_ACCENT = RGBColor(86, 207, 180)
COLOR_WARM = RGBColor(255, 167, 86)
COLOR_TEXT = RGBColor(240, 248, 255)
COLOR_SUBTEXT = RGBColor(190, 213, 232)


def set_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_header(slide, section_no, title, subtitle=""):
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.2), Inches(12.65), Inches(0.95))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_PANEL
    banner.line.fill.background()

    sec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.35), Inches(0.8), Inches(0.55))
    sec_box.fill.solid()
    sec_box.fill.fore_color.rgb = COLOR_ACCENT
    sec_box.line.fill.background()
    sec_tf = sec_box.text_frame
    sec_tf.clear()
    sec_p = sec_tf.paragraphs[0]
    sec_p.text = section_no
    sec_p.font.bold = True
    sec_p.font.size = Pt(20)
    sec_p.font.color.rgb = RGBColor(0, 45, 38)
    sec_p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(1.55), Inches(0.34), Inches(9.8), Inches(0.35))
    ttf = title_box.text_frame
    ttf.clear()
    p = ttf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TEXT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.56), Inches(0.68), Inches(10.4), Inches(0.25))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = COLOR_SUBTEXT


def add_footer(slide, note="PulseGrid Predictive Maintenance Platform"):
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.0), Inches(12.4), Inches(0.3))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{note}  |  Generated {datetime.now().strftime('%d %b %Y')}"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_SUBTEXT


def add_bullets(slide, items, left=0.7, top=1.4, width=5.9, height=5.2, level0_size=20, level1_size=16):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_PANEL
    panel.line.color.rgb = RGBColor(58, 104, 138)

    box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(top + 0.3), Inches(width - 0.55), Inches(height - 0.45))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.color.rgb = COLOR_TEXT if level == 0 else COLOR_SUBTEXT
        p.font.size = Pt(level0_size if level == 0 else level1_size)
        p.space_after = Pt(8 if level == 0 else 4)


def add_metric_cards(slide, cards, start_left=6.95, start_top=1.6):
    x = start_left
    y = start_top
    for i, (title, value, accent) in enumerate(cards):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(14, 38, 58)
        box.line.color.rgb = accent

        t1 = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(2.45), Inches(0.3))
        tf1 = t1.text_frame
        tf1.clear()
        p1 = tf1.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_SUBTEXT

        t2 = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.50), Inches(2.45), Inches(0.55))
        tf2 = t2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = value
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = accent

        x += 2.95
        if (i + 1) % 2 == 0:
            x = start_left
            y += 1.55


def add_simple_table(slide, headers, rows, left=6.9, top=1.6, width=6.0, height=4.6):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(28, 80, 116)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 42, 63)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.color.rgb = COLOR_SUBTEXT


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "00", "PulseGrid", "Predictive Maintenance Platform for Industrial Equipment")
add_bullets(
    slide,
    [
        "End-to-end monitoring and failure-risk prediction system",
        "Supports multi-machine telemetry, alerts, and maintenance workflow",
        "Covers noisy/partial data handling and model monitoring",
        "Prepared by: Daksh Jeena",
    ],
    left=0.7,
    top=1.45,
    width=6.3,
    height=4.9,
    level0_size=20,
)
add_metric_cards(
    slide,
    [
        ("Backend", "Node + Express", COLOR_ACCENT),
        ("Frontend", "React Dashboard", COLOR_WARM),
        ("Database", "JSON Store", COLOR_ACCENT),
        ("Use Case", "Industry 4.0", COLOR_WARM),
    ],
)
add_footer(slide)

# Slide 2: Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "01", "Problem Statement", "Unplanned downtime causes production and revenue loss")
add_bullets(
    slide,
    [
        "Industrial assets fail without enough early warning",
        "Reactive maintenance increases cost and production disruption",
        "Need a real-time predictive maintenance platform that:",
        ("Ingests sensor telemetry (vibration, temperature, current, RPM)", 1),
        ("Predicts failure risk and triggers actionable alerts", 1),
        ("Tracks health trend and recommended maintenance actions", 1),
        ("Works reliably under noisy and missing data", 1),
    ],
)
add_footer(slide)

# Slide 3: Objectives
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "02", "Project Objectives", "Functional and technical goals")
add_bullets(
    slide,
    [
        "Time-series feature extraction and fault-risk modeling",
        "Production-style ingestion and monitoring APIs",
        "Alert generation + maintenance workflow",
        "Model performance and degradation tracking",
        "Evaluation on accuracy, alert lead time, false alarm rate",
    ],
)
add_metric_cards(
    slide,
    [
        ("Accuracy", "Tracked", COLOR_ACCENT),
        ("Lead Time", "Estimated", COLOR_WARM),
        ("False Alarms", "Reduced", COLOR_ACCENT),
        ("Reliability", "Noisy/Partial", COLOR_WARM),
    ],
)
add_footer(slide)

# Slide 4: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "03", "System Architecture", "Data ingestion to insights and alerting")
add_bullets(
    slide,
    [
        "Sensor Stream -> API Ingestion -> Feature & Risk Engine",
        "Risk Engine -> Health Score + Trend + Recommendations",
        "Alert Module -> Open / Resolve maintenance alerts",
        "Dashboard -> fleet view, asset drilldown, monitoring panel",
        "Persistent DB -> assets, readings, alerts, metrics",
    ],
)

a = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(7.0), Inches(1.9), Inches(2.1), Inches(0.8))
a.fill.solid(); a.fill.fore_color.rgb = RGBColor(33, 88, 114); a.line.fill.background(); a.text = "Sensors"
b = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(9.2), Inches(1.9), Inches(2.1), Inches(0.8))
b.fill.solid(); b.fill.fore_color.rgb = RGBColor(43, 110, 131); b.line.fill.background(); b.text = "API"
c = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(7.0), Inches(3.0), Inches(2.1), Inches(0.8))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(58, 131, 142); c.line.fill.background(); c.text = "Risk Engine"
d = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(9.2), Inches(3.0), Inches(2.1), Inches(0.8))
d.fill.solid(); d.fill.fore_color.rgb = RGBColor(68, 152, 147); d.line.fill.background(); d.text = "Alerts"
e = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(8.1), Inches(4.1), Inches(2.1), Inches(0.8))
e.fill.solid(); e.fill.fore_color.rgb = RGBColor(255, 167, 86); e.line.fill.background(); e.text = "Dashboard"
for shp in [a,b,c,d,e]:
    tf = shp.text_frame
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_TEXT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
add_footer(slide)

# Slide 5: Backend/API
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "04", "Backend and API Layer", "Node.js/Express with v2 maintenance endpoints")
add_bullets(
    slide,
    [
        "New routes implemented under /api/v2",
        "Core APIs:",
        ("GET /api/v2/dashboard", 1),
        ("POST /api/v2/readings", 1),
        ("POST /api/v2/readings/simulate", 1),
        ("GET /api/v2/alerts + PATCH alert status", 1),
        "Validation and scoring executed per incoming record",
        "Threshold-driven alerts with maintenance recommendations",
    ],
)
add_simple_table(
    slide,
    ["Endpoint", "Purpose"],
    [
        ["GET /api/v2/health", "Service check"],
        ["GET /api/v2/assets", "List assets"],
        ["POST /api/v2/assets", "Create asset"],
        ["GET /api/v2/readings", "Trend data"],
        ["POST /api/v2/readings", "Ingest telemetry"],
        ["GET /api/v2/alerts", "Open alerts"],
    ],
)
add_footer(slide)

# Slide 6: Database
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "05", "New Database Design", "Persistent local store created for this version")
add_bullets(
    slide,
    [
        "Database file: backend/data/pdm_database.json",
        "Collections (logical sections):",
        ("assets", 1),
        ("sensorReadings", 1),
        ("alerts", 1),
        ("monitoring", 1),
        "Benefits: easy local deployment, reproducible demos, quick prototyping",
        "Can be migrated to MongoDB/PostgreSQL later with same API contract",
    ],
)
add_metric_cards(
    slide,
    [
        ("Storage", "Persistent", COLOR_ACCENT),
        ("Setup", "Zero-config", COLOR_WARM),
        ("Entities", "4 Core", COLOR_ACCENT),
        ("Migration", "DB-ready", COLOR_WARM),
    ],
)
add_footer(slide)

# Slide 7: Frontend
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "06", "Brand-New Frontend", "React dashboard connected to v2 APIs")
add_bullets(
    slide,
    [
        "Fleet KPI cards: assets, alerts, health, monitoring accuracy",
        "Asset panel with risk level chips and health/risk status",
        "Time-series charts for risk/health and sensor trends",
        "Manual ingestion form for real-time reading submission",
        "Asset creation form + simulation trigger",
        "Alert panel with resolve action",
    ],
)
add_simple_table(
    slide,
    ["UI Module", "Connected API"],
    [
        ["Dashboard Summary", "GET /api/v2/dashboard"],
        ["Asset Selector", "GET /api/v2/assets"],
        ["Trend Charts", "GET /api/v2/readings"],
        ["Reading Form", "POST /api/v2/readings"],
        ["Alerts", "GET/PATCH /api/v2/alerts"],
    ],
)
add_footer(slide)

# Slide 8: Analytics
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "07", "Analytics and Monitoring", "What the model and monitoring layer tracks")
add_bullets(
    slide,
    [
        "Risk score (0-1) from sensor-normalized weighted logic",
        "Health score (0-100) with quality penalties",
        "Trend state: rising / stable / improving",
        "Calibration checks: out-of-range sensor detection",
        "Missing-data handling: fallback to last known value",
        "Monitoring: confusion matrix + baseline/current accuracy",
        "Degradation % computed after baseline is established",
    ],
)
add_metric_cards(
    slide,
    [
        ("Risk", "0 - 1", COLOR_WARM),
        ("Health", "0 - 100", COLOR_ACCENT),
        ("Trend", "3 States", COLOR_WARM),
        ("Drift", "% Drop", COLOR_ACCENT),
    ],
)
add_footer(slide)

# Slide 9: Workflow
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "08", "Maintenance Workflow", "From data to action")
add_bullets(
    slide,
    [
        "Step 1: Ingest reading (manual/simulated/real sensor)",
        "Step 2: Score risk and compute health",
        "Step 3: Detect anomalies (missing/calibration)",
        "Step 4: Open alert if threshold crossed",
        "Step 5: Show recommended maintenance action",
        "Step 6: Team resolves alert after intervention",
        "Step 7: Monitoring metrics update from labeled outcomes",
    ],
)
for idx, txt in enumerate([
    "Data In",
    "Risk Score",
    "Alert",
    "Action",
    "Resolve",
]):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0 + (idx%2)*2.8), Inches(1.7 + (idx//2)*1.45), Inches(2.55), Inches(0.95))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(36 + idx*14, 90 + idx*12, 115 + idx*6)
    shp.line.fill.background()
    p = shp.text_frame.paragraphs[0]
    p.text = txt
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT
    p.alignment = PP_ALIGN.CENTER
add_footer(slide)

# Slide 10: Build and Run
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "09", "Build, Run, and Compile", "Execution instructions")
add_bullets(
    slide,
    [
        "Backend start: cd backend && node index.js",
        "Frontend start: cd frontend && npm start",
        "Frontend compile: cd frontend && npm run build",
        "Service check: GET /api/v2/health",
        "Compiled artifacts include frontend/build and docs",
        "Release package created as ZIP for submission/sharing",
    ],
)
add_metric_cards(
    slide,
    [
        ("Backend Port", "8000", COLOR_ACCENT),
        ("Frontend Port", "3000", COLOR_WARM),
        ("Health API", "200 OK", COLOR_ACCENT),
        ("Build", "Successful", COLOR_WARM),
    ],
)
add_footer(slide)

# Slide 11: Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "10", "Expected Results and Marking Points", "Evaluation checklist alignment")
add_bullets(
    slide,
    [
        "Marking Point 1: End-to-end architecture implemented",
        "Marking Point 2: Multi-machine support and ingestion APIs",
        "Marking Point 3: Risk, health, trends, and alert workflow",
        "Marking Point 4: Missing/calibration handling",
        "Marking Point 5: Monitoring and degradation tracking",
        "Marking Point 6: Compiled build + documentation + presentation",
    ],
)
add_simple_table(
    slide,
    ["Evaluation Metric", "Status"],
    [
        ["Accuracy monitoring", "Implemented"],
        ["Alert lead support", "Implemented"],
        ["False alarm tracking", "Implemented"],
        ["Noisy/partial robustness", "Simulation support"],
        ["Production-style flow", "Implemented"],
    ],
)
add_footer(slide)

# Slide 12: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_header(slide, "11", "Conclusion and Next Steps", "Project handoff summary")
add_bullets(
    slide,
    [
        "PulseGrid delivers a complete predictive-maintenance workflow",
        "New frontend + new database + new v2 API are integrated",
        "System is ready for demo, evaluation, and extension",
        "Next steps:",
        ("Connect live MQTT stream", 1),
        ("Swap JSON store with managed DB", 1),
        ("Add advanced ML model retraining pipeline", 1),
        ("Deploy with Docker and CI/CD", 1),
    ],
)
add_metric_cards(
    slide,
    [
        ("Status", "Running", COLOR_ACCENT),
        ("Package", "Prepared", COLOR_WARM),
        ("Docs", "Included", COLOR_ACCENT),
        ("PPT", "Generated", COLOR_WARM),
    ],
)
add_footer(slide, note="Thank you")

prs.save(OUT_FILE)
print(OUT_FILE)
